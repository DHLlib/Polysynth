#!/usr/bin/python3
# -*- coding:utf-8 -*-
"""文件文本提取模块。支持 txt/md/pdf/docx/xlsx/pptx。"""

import os

from backend.core.logger import get_logger

logger = get_logger("file_parser")

_EXTENSION_HANDLERS: dict[str, callable] = {}


def _register(ext: str):
    def decorator(fn):
        _EXTENSION_HANDLERS[ext.lower()] = fn
        return fn
    return decorator


@_register("txt")
@_register("md")
def _parse_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


@_register("pdf")
def _parse_pdf(file_path: str) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        texts = []
        for page in doc:
            texts.append(page.get_text())
        doc.close()
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"PDF parse failed: {e}")
        return ""


@_register("docx")
def _parse_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception as e:
        logger.error(f"DOCX parse failed: {e}")
        return ""


@_register("xlsx")
def _parse_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    lines.append(row_text)
        wb.close()
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"XLSX parse failed: {e}")
        return ""


@_register("pptx")
def _parse_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"PPTX parse failed: {e}")
        return ""


_ALLOWED_EXTENSIONS = {"txt", "md", "pdf", "docx", "xlsx", "pptx"}


def get_file_extension(filename: str) -> str:
    return os.path.splitext(filename)[1].lstrip(".").lower()


def is_allowed_file(filename: str) -> bool:
    return get_file_extension(filename) in _ALLOWED_EXTENSIONS


async def extract_text(file_path: str, filename: str) -> str:
    """从文件中提取纯文本。"""
    ext = get_file_extension(filename)
    if ext not in _ALLOWED_EXTENSIONS:
        logger.warning(f"Unsupported file type: {ext}")
        return ""

    handler = _EXTENSION_HANDLERS.get(ext)
    if not handler:
        return ""

    try:
        text = handler(file_path)
        logger.info(f"Extracted text: {filename}, ext={ext}, len={len(text)}")
        return text
    except Exception as e:
        logger.error(f"Text extraction failed: {filename}, error={e}")
        return ""
